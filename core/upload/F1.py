"""
F1. 조직 데이터 파싱 및 역량 프로필 생성  (개정판)
==================================================

조직소개서(pptx) · 인력정보/특허목록(xlsx)을 받아 표준화된 조직 역량 프로필 JSON을 만든다.

■ 원본에서 고친 것 6건
  ① org_context 플래그가 pptx 문자열 검색에만 의존해 조용히 False가 되었다.
     실측: pptx 파싱 실패 시 inference_speed_ok·has_commercialization_experience가 False.
     그런데 근거는 xlsx에 다 있다 ('30초 내 웹 응답' → 6.팀역량종합,
     '엑스칼리버 국내 최초 상용화' → 1.조직개요).
     → xlsx를 1순위, pptx를 보조 근거로 바꿨다. 각 플래그의 판정 출처를 함께 남긴다.
     ※ has_commercialization_experience가 False면 F2-5에서 진입장벽 F 5점 +
       사업성 F 4점 = 9점이 사라진다.

  ② `from pptx import Presentation`이 모듈 최상단에 있어 python-pptx가 없으면
     모듈 자체가 import 실패했다 ('파일 하나 없어도 안 죽음'이 커버 못 하는 경우).
     → 함수 안으로 옮겨 지연 import 했다.

  ③ capabilities에 보유수준(★)이 없어 F3가 쓸 수 없었다.
     → 6.팀역량종합의 '점수' 컬럼(0~5)을 읽어 level 필드로 넣는다.

  ④ 특허 미반영 역량을 담당영역 키워드 3개로만 찾아 임상검증 등이 누락됐다.
     → 5.개인별역량의 '특허 미반영 보유 역량' 컬럼도 읽는다.

  ⑤ patent_applied_industries가 F3에서 쓸 수 없는 값이었다.
     ['수의영상진단'] 하나만 나오는데, F3-2는 target_industry(후보 산업)와
     비교하므로 절대 매칭되지 않았다.
     → F3-2 개정판에서 해당 로직을 폐기했다. 이 필드는 리포트용으로만 남긴다.

  ⑥ F3가 쓰는 표준역량ID(CAP_*)와 어휘가 갈렸다.
     → standard_capability_levels로 변환해 함께 출력한다.

■ 원본에서 유지한 것 (판단이 맞았음)
  · _find_header_row의 exact match — 부분 일치는 부제목에 낚인다
  · has_relevant_patent·target_industry를 org_context에서 제외 — 후보에 따라 달라지는 값이다
  · 원본 특허 엑셀(5블록, 분류 없음) 대신 분류가 정리된 시트를 쓴 것

실행: python F1.py
"""

import io
import json
import sys
from pathlib import Path

import pandas as pd

BASE = Path(__file__).resolve().parent          # app/core/upload/


def _find_dir(name: str) -> Path:
    """data · samples 폴더를 찾는다.

    ⚠️ BASE.parent/name으로 고정하면 안 된다. F1이 app/upload/에서 app/core/upload/로
       옮겨지면서 그 경로가 app/core/data를 가리키게 됐고, 폴더가 없으니 CSV 로드가
       조용히 건너뛰어졌다(표준역량 검증·특허분류 매핑이 전부 무효화됨).
       에러가 안 나서 알아채기 어렵다 → 상위 폴더를 훑고, 못 찾으면 첫 후보를 돌려준다.
    """
    for p in (BASE.parent / name,           # core/ 옆
              BASE.parent.parent / name,    # app/ 옆  (현재 구조)
              BASE / name):
        if p.is_dir():
            return p
    return BASE.parent / name


DATA_DIR = _find_dir("data")            # 표준역량_정의.csv · 역량어휘_매핑.csv
SAMPLES_DIR = _find_dir("samples")      # 개발용 조직 데이터 (없어도 된다)


# ---------------------------------------------------------------
# 1. capability 정의 — 6.팀역량종합 시트의 10개 카테고리
# ---------------------------------------------------------------
CAPABILITY_DEFINITIONS = [
    {"capability_id": "cap_001", "name": "의료·수의 영상 AI(진단/검출)", "patent_classes": ["A", "C"]},
    {"capability_id": "cap_002", "name": "데이터 품질·증강·합성",       "patent_classes": ["D", "E"]},
    {"capability_id": "cap_003", "name": "정량계측 알고리즘",           "patent_classes": ["B"]},
    {"capability_id": "cap_004", "name": "음성·오디오 AI",             "patent_classes": ["G"]},
    {"capability_id": "cap_005", "name": "멀티모달 학습",               "patent_classes": ["F"]},
    {"capability_id": "cap_006", "name": "위치기반·모빌리티",           "patent_classes": ["I"]},
    {"capability_id": "cap_007", "name": "바이오/오믹스 분석",          "patent_classes": ["H"]},
    # ⚠️ staff_keyword는 좁게 잡는다. '수의'만 쓰면 강태우('산학협력(수의대 컨소시엄)')와
    #    오재혁('수의 문헌 임베딩')까지 걸려 수의 도메인 지식 보유자로 잘못 잡힌다.
    #    실제 수의 도메인 지식 보유자는 노다혜(수의사 면허)뿐이다.
    {"capability_id": "cap_008", "name": "수의 도메인 지식", "level_fallback": 4,
     "patent_classes": [],
     "staff_keyword": ["수의사 면허", "수의학과", "수의영상의학"]},
    {"capability_id": "cap_009", "name": "클라우드 서빙·경량화", "patent_classes": [],
     "staff_keyword": ["MLOps", "경량화", "서빙", "Triton", "저지연"]},
    {"capability_id": "cap_010", "name": "LLM·RAG", "patent_classes": [],
     "staff_keyword": ["LLM", "RAG", "벡터DB", "프롬프트"]},
    # ④ 추가 — 5.개인별역량의 '특허 미반영 보유 역량'에 있으나 원본이 놓친 역량
    #    6.팀역량종합에 없는 카테고리라 ★를 읽을 수 없어 level_fallback으로 명시한다.
    #    '판독'은 임하늘의 'XAI 판독근거 시각화'까지 걸려 제외했다.
    {"capability_id": "cap_011", "name": "임상 검증·판독 품질관리", "level_fallback": 3,
     "patent_classes": [],
     "staff_keyword": ["임상 검증", "임상검증", "임상시험", "Ground Truth"]},
]

# ⑥ F1 capability → F3 표준역량ID
#   F1은 6.팀역량종합의 10개 카테고리를 쓰고, F3는 신사업DB 태그 177개를 커버하는
#   표준역량 32개를 쓴다. 일부는 1:N으로 갈린다.
#   'derived' = 6.팀역량종합의 ★를 그대로 쓴 값 / 'judged' = 세분화 시 판단으로 정한 값
STANDARD_CAP_MAP = {
    "cap_001": [("CAP_VISION", "derived"), ("CAP_DETECT", "derived"),
                ("CAP_CLASSIFY", "judged")],    # A분류를 진단/검출/분류로 세분
    "cap_002": [("CAP_SYNTH", "derived"), ("CAP_VIZ", "judged")],   # E.촬영보조는 별도
    "cap_003": [("CAP_MEASURE", "derived")],
    "cap_004": [("CAP_AUDIO", "derived")],
    "cap_005": [("CAP_MULTIMODAL", "derived")],
    "cap_006": [("CAP_GEO", "derived")],
    "cap_007": [("CAP_BIO", "derived")],
    "cap_008": [("CAP_DOM_VET", "derived")],
    "cap_009": [("CAP_SERVING", "derived"), ("CAP_EDGE", "judged")],  # 엣지는 실적이 30초급
    "cap_010": [("CAP_LLM", "derived")],
    "cap_011": [("CAP_CLINICAL", "judged")],
}
# 세분화 시 낮춰 잡는 값 — 근거가 약한 쪽
JUDGED_LEVEL = {"CAP_CLASSIFY": 4, "CAP_VIZ": 1, "CAP_EDGE": 2, "CAP_CLINICAL": 3}


# ---------------------------------------------------------------
# 2. 전처리
# ---------------------------------------------------------------

def as_bytes(src) -> bytes:
    """경로 문자열 · Path · 업로드 파일 객체를 모두 bytes로 통일한다.

    ⚠️ 왜 필요한가 — Streamlit의 UploadedFile은 읽으면 위치가 끝으로 가서
       두 번째 read()부터 빈 데이터가 된다. F1은 같은 xlsx를 시트별로
       5번(1.조직개요·2.조직원·3.특허목록·5.개인별역량·6.팀역량종합) 읽고,
       각 시트마다 _find_header_row가 한 번 더 읽으므로 총 10회쯤 접근한다.
       bytes로 한 번 받아두고 읽을 때마다 새 BytesIO를 만들어야 한다.
    """
    if isinstance(src, (bytes, bytearray)):
        return bytes(src)
    if hasattr(src, "read"):                     # UploadedFile · BytesIO 등
        if hasattr(src, "seek"):
            src.seek(0)
        data = src.read()
        if hasattr(src, "seek"):
            src.seek(0)
        return data
    return Path(src).read_bytes()


def _stream(src):
    """읽기용 새 스트림. 같은 원본을 몇 번이든 안전하게 다시 읽을 수 있다."""
    return io.BytesIO(as_bytes(src))


def extract_pptx_text(intro_pptx) -> str:
    """조직소개서 pptx의 텍스트(도형+표)를 합쳐 반환. 경로·업로드 객체 모두 받는다.

    ② python-pptx를 함수 안에서 import한다. 최상단에 두면 라이브러리가 없을 때
       모듈 자체가 import 실패해 다른 함수도 못 쓴다.
    """
    from pptx import Presentation   # 지연 import

    prs = Presentation(_stream(intro_pptx))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            lines.append(cell.text.strip())
    return "\n".join(lines)


def _find_header_row(src, sheet_name: str, marker_cell: str) -> int:
    """marker_cell과 '정확히' 일치하는 셀이 있는 행 번호. (부분 일치는 부제목에 낚인다)"""
    raw = pd.read_excel(_stream(src), sheet_name=sheet_name, header=None)
    for i, row in raw.iterrows():
        if row.astype(str).eq(marker_cell).any():
            return i
    raise ValueError(f"'{sheet_name}' 시트에서 '{marker_cell}' 헤더를 찾지 못했다.")


def _parse_sheet(src, sheet_name: str, marker: str, dropna_col: str) -> pd.DataFrame:
    header_row = _find_header_row(src, sheet_name, marker)
    df = pd.read_excel(_stream(src), sheet_name=sheet_name, header=header_row)
    return df.dropna(subset=[dropna_col])


def parse_staff_excel(src) -> pd.DataFrame:
    """2.조직원 — 인력 목록.

    ⚠️ 시트 아래쪽의 각주 행을 걸러낸다. '이름' 컬럼에 값이 있어서
       dropna로는 안 빠지지만 사람이 아니다.
         '* 특허 건수는 중복 제거한 고유 건수 …'
         '* 특허 0건인 5명은 원본 자료에 대표발명자로 미포함 …'
       그대로 두면 인원이 10명 → 12명으로 잡힌다.
    """
    df = _parse_sheet(src, "2.조직원", "이름", "이름")
    name = df["이름"].astype(str).str.strip()
    is_note = name.str.startswith(("*", "※", "-", "#")) | (name.str.len() > 20)
    return df[~is_note].copy()


# ── 특허 분류 추정 규칙 (1-3) ──────────────────────────────────────
#   업로드된 특허 파일에 '분류' 컬럼이 없을 때 발명의 명칭·요약으로 A~I를 추정한다.
#   순서가 중요하다 — 위에서부터 먼저 걸린 분류를 쓴다.
#   정확도: 우리 양식의 35건을 정답으로 두고 측정 34/35 (97%).
#   ⚠️ 추정값이므로 evidence에 inferred=True로 표시한다.
#   순서가 앞선 규칙이 이긴다. 고유 키워드를 먼저, 흔한 단어를 나중에 둔다.
#   ('검출'은 여러 분류에 나오므로 C는 고유 표현만 쓰고 뒤로 보낸다)
PATENT_CLASS_RULES = [
    ("H.바이오마커", ["바이오 마커", "바이오마커"]),
    ("E.촬영 보조", ["자세 교정"]),
    ("F.멀티모달", ["멀티 모달", "멀티모달"]),
    ("G.음성·오디오", ["음악", "가창", "음향", "음성"]),
    ("B.정량 계측", ["심흉비", "심장의 크기", "TPA", "계산하는 장치", "측정하기 위한"]),
    ("D.데이터 품질·증강·합성", ["라벨 오류", "엑스선 영상 생성", "촬영 정보를 처리",
                        "초해상도", "메타 러닝", "데이터 증강", "클래스 불균형",
                        "학습 데이터"]),
    ("C.검출·분할", ["결석", "인스턴스 분할", "객체 탐지", "분할모델"]),
    ("I.위치기반·모빌리티", ["POI", "경로", "차량", "운전", "주행", "위치기반"]),
    ("A.영상 질환 진단", ["질환", "질병", "장기", "진단", "엑스레이", "방사선"]),
]


def infer_patent_class(title: str, summary: str = "") -> str:
    """발명의 명칭·요약에서 특허 분류(A~I)를 추정한다. 못 찾으면 None.

    명칭을 먼저 보고, 명칭으로 안 잡히면 요약을 본다.
    요약을 함께 넣고 한 번에 매칭하면 요약의 부수적 단어가 이긴다
    (예: '라벨 오류 검출' 특허의 요약에 '결석'이 나와 C로 잘못 분류됐다).
    """
    for text in (title, f"{title} {summary}"):
        for cls, keywords in PATENT_CLASS_RULES:
            if any(k in text for k in keywords):
                return cls
    return None


def _find_col(df: pd.DataFrame, *keywords) -> str:
    """컬럼명에 keyword가 들어간 첫 컬럼. 양식이 조금 달라도 찾아낸다."""
    for c in df.columns:
        if any(k in str(c) for k in keywords):
            return c
    return None


def parse_patent_excel(src) -> pd.DataFrame:
    """특허 목록을 읽는다. 두 형태를 모두 받는다.

    ① 우리 양식 — '3.특허목록' 시트, '분류'(A~I)·'발명의 명칭'·'요약' 컬럼
    ② 원본 특허 파일 — 시트 이름 임의, '분류' 컬럼 없음
       (예: digital healthcare팀_특허_보안해제.xlsx → Sheet1에 발명의명칭·요약만)

    ②일 때는 infer_patent_class()로 분류를 추정하고 '분류추정' 컬럼을 True로 둔다.
    ⚠️ 원본 파일은 '순번이 5회 리셋되는 5블록' 구조라 No가 중복된다.
       그래서 No를 그대로 쓰지 않고 행 순서로 다시 번호를 매긴다.
    """
    sheets = list_sheets(src)

    # ① 우리 양식
    if "3.특허목록" in sheets:
        return _parse_sheet(src, "3.특허목록", "분류", "No").assign(분류추정=False)

    # ② 원본 특허 파일 — 헤더 행을 찾아 읽는다
    for sheet in sheets:
        raw = pd.read_excel(_stream(src), sheet_name=sheet, header=None)
        header_row = None
        for i, row in raw.iterrows():
            cells = row.astype(str)
            if cells.str.contains("발명의", na=False).any():
                header_row = i
                break
        if header_row is None:
            continue

        df = pd.read_excel(_stream(src), sheet_name=sheet, header=header_row)
        name_col = _find_col(df, "발명의")
        sum_col = _find_col(df, "요약")
        if name_col is None:
            continue

        df = df.dropna(subset=[name_col])
        df = df.rename(columns={name_col: "발명의 명칭"})
        if sum_col:
            df = df.rename(columns={sum_col: "요약"})
        else:
            df["요약"] = ""

        # ⚠️ 발명의 명칭으로 중복을 제거한다.
        #    원본 파일은 '발명자 5명 기준 총 엔트리 55건 / 고유 35건'(조직개요) 구조라
        #    한 특허가 공동발명자 수만큼 반복된다. 그대로 세면 특허 수가 부풀려진다.
        n_before = len(df)
        df = df.drop_duplicates(subset=["발명의 명칭"], keep="first")
        if len(df) < n_before:
            print(f"  · 특허 중복 제거: {n_before}건 → {len(df)}건 (공동발명 중복)")

        df["분류"] = [infer_patent_class(str(t), str(s)[:200])
                     for t, s in zip(df["발명의 명칭"], df["요약"])]
        df = df[df["분류"].notna()].copy()
        df["No"] = range(1, len(df) + 1)      # 원본 No는 5블록으로 중복돼 못 쓴다
        df["분류추정"] = True
        return df[["No", "분류", "발명의 명칭", "요약", "분류추정"]]

    raise ValueError(f"특허 목록을 읽을 수 없다. 시트: {sheets}")


def parse_team_capability_excel(src) -> pd.DataFrame:
    """③ 6.팀역량종합 — 역량 카테고리별 수준(★)·점수·특허건수·근거."""
    return _parse_sheet(src, "6.팀역량종합", "역량 카테고리", "역량 카테고리")


def parse_individual_capability_excel(src) -> pd.DataFrame:
    """④ 5.개인별역량 — '특허가 입증하는 역량' / '특허 미반영 보유 역량'."""
    return _parse_sheet(src, "5.개인별역량", "이름", "이름")


def list_sheets(src) -> list:
    """xlsx에 들어 있는 시트 이름. 업로드 파일이 우리 양식인지 판별할 때 쓴다."""
    return pd.ExcelFile(_stream(src)).sheet_names


def parse_org_overview_excel(src) -> dict:
    """1.조직개요 — 구분/내용 2열 key-value."""
    raw = pd.read_excel(_stream(src), sheet_name="1.조직개요", header=None)
    out = {}
    for _, row in raw.iterrows():
        k, v = str(row.iloc[0]).strip(), str(row.iloc[1]).strip()
        if k and k != "nan" and v != "nan":
            out[k] = v
    return out


# ---------------------------------------------------------------
# 3. 역량 매핑
# ---------------------------------------------------------------

def map_patents_to_capabilities(patent_df: pd.DataFrame) -> tuple:
    """특허 → (capability_id별 evidence_id, evidence_id별 상세)"""
    cap_evidence = {c["capability_id"]: [] for c in CAPABILITY_DEFINITIONS}
    evidence_detail = {}

    class_to_cap = {cls: cap["capability_id"]
                    for cap in CAPABILITY_DEFINITIONS for cls in cap["patent_classes"]}

    for _, row in patent_df.iterrows():
        classification = str(row.get("분류", "")).strip()
        cap_id = class_to_cap.get(classification[0] if classification else None)
        if cap_id is None:
            continue
        no = int(float(row.get("No")))
        eid = f"ev_patent_{no:03d}"
        cap_evidence[cap_id].append(eid)
        evidence_detail[eid] = {
            "type": "patent", "patent_no": no,
            "title": str(row.get("발명의 명칭", "")).strip(),
            "classification": classification,
            # 분류가 우리 양식에 있던 값인지, 명칭·요약으로 추정한 값인지 남긴다
            "class_inferred": bool(row.get("분류추정", False)),
            "summary": str(row.get("요약", ""))[:300],
        }
    return cap_evidence, evidence_detail


def enrich_from_staff(staff_df, indiv_df, cap_evidence: dict, evidence_detail: dict) -> None:
    """특허 0건 역량을 인력정보에서 보강한다.

    ④ 원본은 2.조직원의 '담당 영역'만 봤다. 5.개인별역량의
       '특허 미반영 보유 역량' 컬럼에 임상검증·XAI·ONNX 등이 더 있어 함께 읽는다.
    """
    keyword_by_cap = {c["capability_id"]: c.get("staff_keyword", [])
                      for c in CAPABILITY_DEFINITIONS}

    # 이름 → 검색 대상 텍스트
    texts = {}
    for _, r in staff_df.iterrows():
        name = str(r.get("이름", "")).strip()
        if name and name != "nan":
            texts[name] = [str(r.get("담당 영역", "")), str(r.get("경력", "")),
                           str(r.get("학력", ""))]
    if indiv_df is not None:
        cols = list(indiv_df.columns)
        unref = next((c for c in cols if "미반영" in str(c)), None)
        proven = next((c for c in cols if "입증" in str(c)), None)
        for _, r in indiv_df.iterrows():
            name = str(r.get("이름", "")).strip()
            if name in texts:
                for c in (unref, proven):
                    if c:
                        texts[name].append(str(r.get(c, "")))

    for name, parts in texts.items():
        blob = " ".join(parts)
        for cap_id, keywords in keyword_by_cap.items():
            if not keywords:
                continue
            hit = next((kw for kw in keywords if kw in blob), None)
            if hit is None:
                continue
            eid = f"ev_staff_{name}"
            if eid not in cap_evidence[cap_id]:
                cap_evidence[cap_id].append(eid)
            if eid not in evidence_detail:
                evidence_detail[eid] = {"type": "staff", "name": name,
                                        "matched_keyword": [], "text": blob[:300]}
            evidence_detail[eid]["matched_keyword"].append(f"{cap_id}:{hit}")


def _norm_cap(s) -> str:
    """카테고리명 비교용 정규화 — 공백·괄호를 지우고 소문자로."""
    return str(s).replace(" ", "").replace("(", "").replace(")", "").lower()


# 역량이 아닌 집계·주석 행을 걸러낸다
_NOT_A_CAPABILITY = ("합계", "소계", "총계", "평균", "갭 신호", "[")


def read_capability_rows(team_df) -> list:
    """6.팀역량종합에서 (카테고리명, 점수, 근거) 행을 뽑는다.

    read_capability_levels(매칭)와 unmatched_capability_rows(미매칭)가 같은 목록을
    봐야 한다. 따로 읽으면 한쪽만 고쳐져서 매칭·미매칭 합이 안 맞는다.
    """
    if team_df is None:
        return []
    cols = list(team_df.columns)
    name_col = cols[0]
    score_col = next((c for c in cols if "점수" in str(c)), None)
    note_col = next((c for c in cols if "근거" in str(c)), None)
    if score_col is None:
        return []

    rows = []
    for _, r in team_df.iterrows():
        v = pd.to_numeric(r[score_col], errors="coerce")
        if pd.isna(v):
            continue
        name = str(r[name_col]).strip()
        if not name or name == "nan":
            continue
        if any(t in name for t in _NOT_A_CAPABILITY) or len(name) > 40:
            continue        # 합계행·주석행은 역량이 아니다
        rows.append({"name": name, "level": int(v),
                     "note": str(r[note_col]) if note_col else ""})
    return rows


def match_capability_name(name: str) -> str:
    """카테고리명 → F1 capability_id. 못 찾으면 None.

    완전일치 → 부분포함 순. 의미가 아니라 글자를 비교하므로 '영상진단 AI'처럼
    어순이 다르면 못 잡는다. 그 경우는 LLM 폴백이 32개 표준역량 축에서 찾는다.
    """
    k = _norm_cap(name)
    for cap in CAPABILITY_DEFINITIONS:
        if k == _norm_cap(cap["name"]):
            return cap["capability_id"]
    for cap in CAPABILITY_DEFINITIONS:
        key = _norm_cap(cap["name"])
        if key in k or k in key:
            return cap["capability_id"]
    return None


def unmatched_capability_rows(team_df) -> list:
    """11개 카테고리에 매칭되지 않은 시트 행. LLM 폴백의 입력이 된다."""
    return [r for r in read_capability_rows(team_df)
            if match_capability_name(r["name"]) is None]


def read_capability_levels(team_df, evidence_df=None) -> dict:
    """③ 6.팀역량종합의 '점수'(0~5)를 capability_id별로 읽는다.

    시트의 '역량 카테고리' 표기와 CAPABILITY_DEFINITIONS의 name이 완전히 같지 않아
    (예: '정량 계측 알고리즘' vs '정량계측 알고리즘') 공백을 제거하고 비교한다.

    Args:
        team_df    : 6.팀역량종합
        evidence_df: level_fallback의 근거를 찾을 시트들(2.조직원·5.개인별역량).
                     비우면 fallback을 적용하지 않는다.

    ⚠️ level_fallback을 무조건 주면 안 된다. 시트에 ★가 없는 역량(cap_008 수의
       도메인 지식 ★4 · cap_011 임상검증 ★3)에 상수를 넣는 장치인데, 조건 없이
       주면 업로드 내용과 무관하게 항상 부여된다. 수의와 무관한 조직이 자기 파일을
       올려도 '수의 도메인 지식 ★4'를 물려받아, 다른 조직으로 테스트할 때
       실행가능성(C 도메인 전문성 5점)이 실제보다 높게 나온다.
       → staff_keyword가 인력정보에서 실제로 잡힐 때만 인정한다.
    """
    # fallback 근거를 찾을 텍스트 (인력정보·개인별역량 전체)
    evidence_text = " ".join(
        "" if df is None else " ".join(df.astype(str).values.ravel())
        for df in (evidence_df or []))

    matched = {}
    for row in read_capability_rows(team_df):
        cap_id = match_capability_name(row["name"])
        if cap_id and cap_id not in matched:
            matched[cap_id] = row["level"]

    out = dict(matched)
    for cap in CAPABILITY_DEFINITIONS:
        if cap["capability_id"] in out or cap.get("level_fallback") is None:
            continue
        keywords = cap.get("staff_keyword") or []
        if keywords and any(kw in evidence_text for kw in keywords):
            out[cap["capability_id"]] = cap["level_fallback"]
    return out


# ---------------------------------------------------------------
# 3-b. LLM 폴백 — 11개 카테고리에 없는 역량을 표준역량 32개 축에 직접 매핑
#
# ■ 왜 필요한가
#   F1의 카테고리 11개는 코드에 박혀 있고 글자 비교로만 매칭된다. 그래서
#   '영상진단 AI'처럼 어순이 다르거나 '진동·소음 이상탐지'처럼 다른 도메인 역량은
#   조용히 버려진다. 표준역량은 32개인데 11개 축을 거치면 15개만 도달 가능해서,
#   신사업 DB 50건 중 30건이 요구하는 역량을 어떤 파일을 올려도 보유할 수 없다.
#
# ■ 안전장치 3개
#   ① 환각 차단 — 응답을 표준역량ID enum으로 제약하고, 받은 뒤 화이트리스트로
#      한 번 더 검증한다(F4-1·F4-3과 같은 이중 방어).
#      보유수준(★)은 LLM에 묻지 않는다. 시트의 '점수' 값을 그대로 쓴다.
#      수준을 LLM이 만들면 없는 역량이 점수로 둔갑한다.
#   ② 재현성 — temperature=0 + 같은 입력은 캐시에서 돌려준다.
#   ③ 비용 — 문자열 매칭이 성공한 항목은 LLM에 보내지 않는다. 전부 매칭되면
#      (샘플 조직 = 10/10) 호출 자체가 없다.
#
# ■ 꺼지는 조건 — 키가 없거나 openai 미설치거나 호출 실패면 폴백을 건너뛴다.
#   그때는 이 기능이 없던 것과 완전히 동일하게 동작한다(미매칭 행은 버려짐).
# ---------------------------------------------------------------
LLM_MAP_MODEL_ENV = "F1_OPENAI_MODEL"
LLM_MAP_MODEL_DEFAULT = "gpt-5-mini"
_llm_map_cache = {}


def load_standard_capability_choices() -> list:
    """표준역량_정의.csv에서 (ID, 이름, 설명) — LLM에게 보여줄 선택지."""
    path = DATA_DIR / "표준역량_정의.csv"
    if not path.exists():
        return []
    df = pd.read_csv(path)
    desc_col = "LLM용설명" if "LLM용설명" in df.columns else None
    return [{"id": str(r["표준역량ID"]), "name": str(r["표준역량명"]),
             "desc": str(r[desc_col]) if desc_col and pd.notna(r[desc_col]) else ""}
            for _, r in df.iterrows()]


def _llm_map_unmatched(rows: list, choices: list) -> dict:
    """미매칭 카테고리명 → 표준역량ID 리스트. 실패하면 {}."""
    if not rows or not choices:
        return {}

    key = tuple(sorted(r["name"] for r in rows))
    if key in _llm_map_cache:                      # ② 재현성·비용
        return _llm_map_cache[key]

    import os
    if not os.environ.get("OPENAI_API_KEY"):
        # Streamlit Cloud 등 배포 환경에서는 secrets.toml에만 키가 있고
        # os.environ에는 없을 수 있으므로 st.secrets에서 찾아 os.environ에 채워둔다
        # (아래 OpenAI() 생성자가 os.environ만 보기 때문).
        try:
            import streamlit as st

            secret_key = st.secrets.get("OPENAI_API_KEY")
        except Exception:
            secret_key = None
        if secret_key:
            os.environ["OPENAI_API_KEY"] = secret_key
    if not os.environ.get("OPENAI_API_KEY"):
        return {}
    try:
        from openai import OpenAI
    except ImportError:
        return {}

    valid = {c["id"] for c in choices}
    menu = "\n".join(f"- {c['id']}: {c['name']}"
                     + (f" — {c['desc'][:80]}" if c["desc"] else "")
                     for c in choices)
    asked = "\n".join(f"- {r['name']}" + (f" (근거: {r['note'][:60]})" if r["note"] else "")
                      for r in rows)

    schema = {
        "type": "object",
        "properties": {
            "mappings": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "category": {"type": "string",
                                     "enum": [r["name"] for r in rows]},
                        # ① enum 제약 — DB에 없는 ID를 만들 수 없다
                        "standard_capability_ids": {
                            "type": "array",
                            "items": {"type": "string", "enum": sorted(valid)},
                        },
                        "reason": {"type": "string"},
                    },
                    "required": ["category", "standard_capability_ids", "reason"],
                    "additionalProperties": False,
                },
            }
        },
        "required": ["mappings"],
        "additionalProperties": False,
    }

    system = (
        "너는 조직 역량 어휘를 표준 분류에 맞추는 매핑 엔진이다.\n"
        "사용자가 준 '조직이 신고한 역량 이름'을 아래 표준역량 목록에서 "
        "의미가 같은 것으로 연결한다.\n\n"
        "규칙\n"
        "- 목록에 있는 표준역량ID만 사용한다. 새 ID를 만들지 않는다.\n"
        "- 의미가 실제로 겹칠 때만 연결한다. 애매하면 빈 배열로 둔다.\n"
        "- 한 역량이 여러 표준역량으로 갈릴 수 있다(예: 영상 AI → 진단 + 검출).\n"
        "- 억지로 채우지 않는다. 관련 없으면 연결하지 않는 것이 맞다.\n"
        "- 보유 수준·점수는 판단하지 않는다. 이름 매핑만 한다.\n\n"
        f"[표준역량 목록]\n{menu}"
    )
    payload = dict(
        model=os.environ.get(LLM_MAP_MODEL_ENV, LLM_MAP_MODEL_DEFAULT),
        messages=[{"role": "system", "content": system},
                  {"role": "user", "content": f"[조직이 신고한 역량]\n{asked}"}],
        response_format={"type": "json_schema",
                         "json_schema": {"name": "capability_mapping",
                                         "schema": schema, "strict": True}},
    )

    try:
        client = OpenAI(timeout=60.0, max_retries=2)
        try:
            resp = client.chat.completions.create(temperature=0, **payload)
        except Exception as e:
            if "temperature" in str(e):      # 일부 모델은 temperature를 거부한다
                resp = client.chat.completions.create(**payload)
            else:
                raise
        data = json.loads(resp.choices[0].message.content)
    except Exception as e:
        print(f"⚠️ 역량 매핑 LLM 폴백 실패 — 미매칭 역량은 제외됩니다: {e}")
        return {}

    out = {}
    for m in data.get("mappings", []):
        cat = str(m.get("category", "")).strip()
        # ① 후처리 화이트리스트 — enum을 통과했어도 한 번 더 확인
        ids = [i for i in m.get("standard_capability_ids", []) if i in valid]
        if cat and ids:
            out[cat] = {"ids": ids, "reason": str(m.get("reason", ""))[:200]}
    _llm_map_cache[key] = out
    return out


def _class_letter(v) -> str:
    """'A.영상 질환 진단' → 'A'. 분류 표기가 바뀌어도 앞 글자만 본다."""
    s = str(v).strip()
    return s[0] if s and s[0].isalpha() else ""


def count_patents_by_class(patent_df) -> dict:
    """분류(A~I)별 특허 건수. 표준역량별 특허건수의 재료다.

    표준역량_정의.csv의 '특허건수'는 특정 조직 스냅샷이므로 그대로 쓰면 누가
    업로드해도 같은 값이 된다. 업로드한 특허를 분류별로 세서 대체한다.
    (실측: 샘플 데이터 35건을 세면 CSV의 특허건수와 전 항목 일치)
    """
    if patent_df is None or len(patent_df) == 0 or "분류" not in patent_df.columns:
        return {}
    out = {}
    for v in patent_df["분류"]:
        letter = _class_letter(v)
        if letter:
            out[letter] = out.get(letter, 0) + 1
    return out


def inventors_by_class(patent_df) -> dict:
    """분류(A~I)별 발명자 이름. 표준역량별 '전담인력' 판정의 재료다.

    특허 시트는 발명자 컬럼에 'O' 표시로 참여를 나타낸다(이름 컬럼 = 조직원 이름).
    """
    if patent_df is None or len(patent_df) == 0 or "분류" not in patent_df.columns:
        return {}
    skip = {"No", "분류", "발명의 명칭", "발명자 수", "요약", "분류추정"}
    name_cols = [c for c in patent_df.columns if str(c) not in skip]

    out = {}
    for _, row in patent_df.iterrows():
        letter = _class_letter(row.get("분류"))
        if not letter:
            continue
        names = out.setdefault(letter, [])
        for c in name_cols:
            if pd.notna(row.get(c)) and str(row.get(c)).strip():
                if str(c) not in names:
                    names.append(str(c))
    return out


# 6.팀역량종합 '근거' 텍스트 → 로드맵 연계 등급.
#   ⚠️ '확장'·'보조'는 업로드 양식에 신호가 없어 유도할 수 없다(핵심으로 잡힌다).
#      샘플 데이터 기준 12개 중 10개가 CSV와 일치하고, CAP_VIZ(보조)·
#      CAP_MULTIMODAL(확장)만 핵심으로 올라간다. 배점 영향은 B로드맵연계 4점뿐이다.
ROADMAP_RULES = [
    ("미연계", ["로드맵 미활용", "미활용", "레거시"]),
    ("개발 단계", ["개발 단계", "착수", "초기 개발"]),
]


def read_roadmap_notes(team_df) -> dict:
    """6.팀역량종합의 '근거' 문장을 capability_id별로 읽는다."""
    if team_df is None:
        return {}
    cols = list(team_df.columns)
    name_col = cols[0]
    note_col = next((c for c in cols if "근거" in str(c)), None)
    if note_col is None:
        return {}

    def norm(s):
        return str(s).replace(" ", "").replace("(", "").replace(")", "").lower()

    sheet = {norm(r[name_col]): str(r[note_col]) for _, r in team_df.iterrows()}
    out = {}
    for cap in CAPABILITY_DEFINITIONS:
        key = norm(cap["name"])
        note = sheet.get(key)
        if note is None:
            note = next((v for k, v in sheet.items() if key in k or k in key), None)
        if note and note != "nan":
            out[cap["capability_id"]] = note
    return out


def _roadmap_grade(note: str, patent_count: int):
    """근거 문장 + 특허 보유로 로드맵 연계 등급을 정한다.

    특허가 0건이면 None. CSV도 특허 0건 역량은 로드맵연계를 비워두며,
    F3-2의 score_roadmap_link가 빈 값을 가중치 0으로 처리한다.
    """
    if not patent_count:
        return None
    for grade, keywords in ROADMAP_RULES:
        if note and any(kw in note for kw in keywords):
            return grade
    return "핵심"


def to_standard_capabilities(cap_levels: dict, patent_df=None, team_df=None,
                             cap_evidence: dict = None,
                             evidence_detail: dict = None,
                             use_llm: bool = True) -> dict:
    """⑥ F1 capability → F3 표준역량ID(CAP_*)별 조직 사실.

    F3-2가 채점에 쓰는 조직 정보를 전부 업로드 데이터에서 만들어 넘긴다.
      level        6.팀역량종합 ★ (judged 분기는 JUDGED_LEVEL)
      patent_count 업로드 특허를 분류(A~I)별로 센 값
      staff        그 분류 특허의 발명자 + staff_keyword로 잡은 인력
      roadmap      6.팀역량종합 '근거' 문장에서 유도한 연계 등급

    표준역량_정의.csv는 '어느 표준역량이 어느 특허분류인가' 같은 정의(분류 체계)로만
    쓴다. 조직마다 달라지는 값(수준·특허수·인력·로드맵)은 업로드에서 나온다.
    표준역량_정의.csv에 없는 ID가 나오면 경고한다 (어휘가 또 갈리는 것을 막는다).

    Args:
        use_llm: 11개 카테고리에 매칭 안 된 역량을 LLM으로 표준역량 축에 매핑한다.
                 매칭이 전부 성공하면 호출이 없다(샘플 조직 = 10/10 → 비용 0).
                 키가 없거나 실패하면 조용히 건너뛴다(= 이 기능이 없던 동작).
    """
    known, class_of = None, {}
    path = DATA_DIR / "표준역량_정의.csv"
    if path.exists():
        std = pd.read_csv(path)
        known = set(std["표준역량ID"])
        if "특허분류" in std.columns:
            class_of = {str(r["표준역량ID"]): _class_letter(r["특허분류"])
                        for _, r in std.iterrows() if pd.notna(r["특허분류"])}

    class_counts = count_patents_by_class(patent_df)
    class_inventors = inventors_by_class(patent_df)
    notes = read_roadmap_notes(team_df)

    # staff_keyword로 잡힌 인력 (특허가 없는 역량 — 클라우드 서빙·LLM·임상검증 등)
    kw_staff = {}
    for cap_id, eids in (cap_evidence or {}).items():
        for eid in eids:
            ev = (evidence_detail or {}).get(eid, {})
            if ev.get("type") == "staff" and ev.get("name"):
                kw_staff.setdefault(cap_id, []).append(ev["name"])

    # ── ⓐ 문자열 매칭으로 잡힌 역량 (기존 경로 — 비용 0) ──
    plan = []       # (std_id, level, source, from_label, note, cap_id)
    for cap_id, targets in STANDARD_CAP_MAP.items():
        base = cap_levels.get(cap_id)
        for std_id, kind in targets:
            if known is not None and std_id not in known:
                print(f"⚠️ 표준역량_정의.csv에 없는 ID: {std_id}")
            if kind == "judged":
                # ⚠️ 부모 역량이 없으면 세분화할 것도 없다. base를 안 보고
                #    JUDGED_LEVEL을 그대로 쓰면, 그 역량을 신고하지 않은 조직에도
                #    CAP_CLINICAL ★3·CAP_VIZ ★1이 근거 없이 붙는다(실측).
                if base is None:
                    continue
                lv, src = JUDGED_LEVEL.get(std_id, base), "판단(세분화)"
            else:
                lv, src = base, "6.팀역량종합 점수"
            if lv is not None:
                plan.append((std_id, lv, src, cap_id, notes.get(cap_id, ""), cap_id))

    # ── ⓑ 매칭 실패한 역량만 LLM에게 물어본다 (미매칭 0건이면 호출 없음) ──
    if use_llm:
        unmatched = unmatched_capability_rows(team_df)
        if unmatched:
            mapping = _llm_map_unmatched(unmatched, load_standard_capability_choices())
            taken = {p[0] for p in plan}
            for row in unmatched:
                hit = mapping.get(row["name"])
                if not hit:
                    continue
                for std_id in hit["ids"]:
                    if std_id in taken:
                        continue    # 문자열 매칭 결과를 LLM이 덮어쓰지 못하게 한다
                    taken.add(std_id)
                    plan.append((std_id, row["level"], "LLM 매핑",
                                 row["name"], row["note"], None))

    # ── ⓒ 특허수·인력·로드맵을 붙인다 (두 경로 동일하게) ──
    out = {}
    for std_id, lv, src, from_label, note, cap_id in plan:
        letter = class_of.get(std_id, "")
        n_patent = class_counts.get(letter, 0) if letter else 0
        staff = list(class_inventors.get(letter, [])) if letter else []
        for nm in kw_staff.get(cap_id, []) if cap_id else []:
            if nm not in staff:
                staff.append(nm)

        out[std_id] = {
            "level": int(lv), "source": src, "from": from_label,
            "patent_class": letter or None,
            "patent_count": n_patent,
            "staff": staff,
            "roadmap": _roadmap_grade(note, n_patent),
        }
    return out


# ---------------------------------------------------------------
# 4. org_context
# ---------------------------------------------------------------

def build_org_context(overview: dict, staff_df, team_df, indiv_df,
                      intro_text: str, patent_df) -> dict:
    """① 조직 플래그를 xlsx 1순위 · pptx 보조로 판정하고 출처를 남긴다.

    원본은 pptx 문자열 검색만 썼고, 파싱이 실패하면 조용히 False가 되었다.
    근거는 xlsx에 다 있다:
      '30초 내 웹 응답'          → 6.팀역량종합
      '엑스칼리버 국내 최초 상용화' → 1.조직개요 대표 성과
      '수의사 면허'              → 2.조직원 학력
      '산학협력(수의대 컨소시엄)'   → 5.개인별역량 특허 미반영 역량
    """
    def blob(df):
        return "" if df is None else " ".join(df.astype(str).values.ravel())

    xlsx_text = " ".join([" ".join(f"{k} {v}" for k, v in overview.items()),
                          blob(staff_df), blob(team_df), blob(indiv_df)])

    def judge(keywords):
        """(값, 출처) — xlsx 우선, 없으면 pptx."""
        for kw in keywords:
            if kw in xlsx_text:
                return True, f"xlsx('{kw}')"
        for kw in keywords:
            if kw in intro_text:
                return True, f"pptx('{kw}')"
        return False, "근거 없음"

    expert, src_expert = judge(["수의사 면허", "수의사"])
    speed, src_speed = judge(["30초", "저지연", "실시간 응답"])
    partner, src_partner = judge(["산학", "컨소시엄", "동물병원"])
    comm, src_comm = judge(["상용화", "엑스칼리버", "X Caliber"])

    n_patent = 0 if patent_df is None else len(patent_df)

    return {
        "has_domain_expert": expert,
        "inference_speed_ok": speed,
        "has_partner_network": partner,
        "has_commercialization_experience": comm,
        "patent_count": n_patent,                     # ③ F3-2가 읽는 키
        "_sources": {"has_domain_expert": src_expert,
                     "inference_speed_ok": src_speed,
                     "has_partner_network": src_partner,
                     "has_commercialization_experience": src_comm},
    }


def detect_ip_gap(capabilities: list, evidence_detail: dict) -> list:
    """멘토님 신호 02 — 역량은 보유했으나 관련 특허 0건 (IP 공백).

    '30초 내 웹 응답, 무서버 병원 배포를 실제 구현했으나 관련 특허 0건.
     보유 역량 대비 IP 공백(방어 취약점)을 탐지하는 케이스.'
    """
    out = []
    for c in capabilities:
        if not c["evidence_ids"]:
            continue
        has_patent = any(evidence_detail[e]["type"] == "patent" for e in c["evidence_ids"])
        if not has_patent:
            out.append({"capability_id": c["capability_id"], "name": c["name"],
                        "level": c.get("level"),
                        "evidence": [evidence_detail[e].get("name", e)
                                     for e in c["evidence_ids"]]})
    return out


def summarize_staff(staff_df) -> list:
    """2-1 미리보기용 인력 목록. 컬럼명이 조금 달라도 찾아낸다."""
    if staff_df is None or staff_df.empty:
        return []
    cols = list(staff_df.columns)

    def col(*keys):
        return next((c for c in cols if any(k in str(c) for k in keys)), None)

    c_name, c_rank = col("이름"), col("직급")
    c_role, c_edu = col("담당"), col("학력")
    c_career, c_pat = col("경력"), col("특허")

    out = []
    for _, r in staff_df.iterrows():
        def get(c):
            return "" if c is None or pd.isna(r.get(c)) else str(r[c]).strip()
        out.append({
            "이름": get(c_name), "직급": get(c_rank), "담당": get(c_role),
            "학력": get(c_edu), "경력": get(c_career),
            "특허": int(pd.to_numeric(r.get(c_pat), errors="coerce") or 0) if c_pat else 0,
        })
    return out


def detect_patent_concentration(staff_df) -> dict:
    """멘토님 신호 03 — 특허가 상위 인력에 편중된 정도."""
    if staff_df is None or staff_df.empty:
        return {}
    col = next((c for c in staff_df.columns if "특허" in str(c)), None)
    if col is None:
        return {}
    pc = pd.to_numeric(staff_df[col], errors="coerce").dropna()
    if pc.empty or pc.sum() == 0:
        return {}
    names = staff_df.loc[pc.index, "이름"].astype(str)
    return {
        "total": int(pc.sum()),
        "headcount": int(len(pc)),
        "top3_share": round(pc.nlargest(3).sum() / pc.sum() * 100, 1),
        "zero_patent_staff": sorted(names[pc == 0].tolist()),
    }


# ---------------------------------------------------------------
# 5. 메인
# ---------------------------------------------------------------

def build_organization_profile(files: dict) -> dict:
    """조직 데이터를 파싱해 역량 프로필을 만든다.

    files의 값은 경로 문자열 · Path · 업로드 파일 객체 아무거나 된다.
      {"intro_pptx": ..., "staff_excel": ..., "patent_excel": ...(생략 가능)}

    반환에 parse_status가 포함된다 (기능명세 1-5 업로드 진행 상태 표시용).
    """
    status = {}     # 1-5  단계별 파싱 결과

    def safe(fn, label, default=None):
        try:
            out = fn()
            status[label] = {"ok": True, "detail": ""}
            return out
        except Exception as e:
            status[label] = {"ok": False, "detail": str(e)[:120]}
            print(f"⚠️ {label} 실패, 건너뜀: {e}")
            return default

    intro_src = files.get("intro_pptx")
    staff_xl = files.get("staff_excel")
    patent_xl = files.get("patent_excel") or staff_xl   # 없으면 인력정보 파일에서 찾는다

    intro_text = safe(lambda: extract_pptx_text(intro_src),
                      "조직소개서", "") if intro_src else ""
    if not intro_src:
        status["조직소개서"] = {"ok": False, "detail": "업로드되지 않음"}

    if staff_xl is None:
        status["인력정보"] = {"ok": False, "detail": "업로드되지 않음"}
        staff_df = team_df = indiv_df = None
        overview = {}
    else:
        staff_df = safe(lambda: parse_staff_excel(staff_xl), "인력정보")
        team_df = safe(lambda: parse_team_capability_excel(staff_xl), "팀역량종합")
        indiv_df = safe(lambda: parse_individual_capability_excel(staff_xl), "개인별역량")
        overview = safe(lambda: parse_org_overview_excel(staff_xl), "조직개요", {})

    if patent_xl is None:
        status["특허목록"] = {"ok": False, "detail": "업로드되지 않음"}
        patent_df = None
    else:
        patent_df = safe(lambda: parse_patent_excel(patent_xl), "특허목록")

    if staff_df is None:
        staff_df = pd.DataFrame(columns=["이름", "담당 영역", "학력", "경력"])
    if patent_df is None:
        patent_df = pd.DataFrame(columns=["No", "분류", "발명의 명칭", "요약"])

    cap_evidence, evidence_detail = map_patents_to_capabilities(patent_df)
    enrich_from_staff(staff_df, indiv_df, cap_evidence, evidence_detail)
    cap_levels = read_capability_levels(team_df, [staff_df, indiv_df])

    capabilities = [{
        "capability_id": c["capability_id"],
        "name": c["name"],
        "level": cap_levels.get(c["capability_id"]),      # ③ 0~5, 없으면 None
        "evidence_ids": cap_evidence.get(c["capability_id"], []),
    } for c in CAPABILITY_DEFINITIONS]

    return {
        "organization_id": "org_001",
        # 2-1 미리보기용 — 원본에서 읽은 조직 개요와 인력 목록.
        #     원래는 org_context 판정에만 쓰고 버렸는데, 사용자가 '무엇을 읽었는지'
        #     확인할 수 있어야 하므로 프로필에 담는다.
        "organization": {str(k): str(v) for k, v in overview.items()
                         if not str(k).startswith("[")},
        "staff": summarize_staff(staff_df),
        "capabilities": capabilities,
        "evidence_ids": sorted(evidence_detail.keys()),
        "evidence": evidence_detail,
        # ⑥ F3 연결용 — 표준역량별 수준·특허수·인력·로드맵을 모두 업로드에서 만든다
        "standard_capability_levels": to_standard_capabilities(
            cap_levels, patent_df, team_df, cap_evidence, evidence_detail),
        # 분류(A~I)별 특허 건수 — 표준역량 단위로 쪼개기 전 원본 집계 (검증·표시용)
        "patent_class_counts": count_patents_by_class(patent_df),
        "org_context": build_org_context(overview, staff_df, team_df, indiv_df,
                                        intro_text, patent_df),
        "signals": {                                    # 멘토님 '잡아내야 할 신호'
            "ip_gap": detect_ip_gap(capabilities, evidence_detail),
            "patent_concentration": detect_patent_concentration(staff_df),
        },
        # 1-4  특허 보유 여부 자동 판별 — 파일 업로드 여부가 아니라
        #      실제로 분류(A~I)가 붙은 특허를 읽었는지로 판정한다.
        "has_patent_data": len(patent_df) > 0,
        # 1-5  업로드 진행 상태 표시
        "parse_status": status,
    }


if __name__ == "__main__":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except AttributeError:
        pass

    # 개발용 — samples/ 폴더가 있으면 거기서 읽는다. 실제로는 업로드로 받는다.
    ORG = SAMPLES_DIR if SAMPLES_DIR.exists() else Path.home() / "Desktop" / "A19" / "조직데이터"
    files = {
        "intro_pptx": str(next(ORG.glob("*.pptx"))),
        "staff_excel": str(next(ORG.glob("*org_data*.xlsx"))),
    }
    p = build_organization_profile(files)

    print(f"\ncapabilities {len(p['capabilities'])}개 · evidence {len(p['evidence_ids'])}건\n")
    print(f"  {'ID':<9}{'역량명':<26}{'★':>3}{'evidence':>10}")
    print("  " + "-" * 50)
    for c in p["capabilities"]:
        lv = "—" if c["level"] is None else str(c["level"])
        print(f"  {c['capability_id']:<9}{c['name'][:24]:<26}{lv:>3}{len(c['evidence_ids']):>8}건")

    print(f"\n표준역량 변환 ({len(p['standard_capability_levels'])}개) — F3 연결용")
    for k, v in sorted(p["standard_capability_levels"].items()):
        print(f"  {k:<18}★{v['level']}  {v['source']:<16}← {v['from']}")

    print("\norg_context")
    src = p["org_context"]["_sources"]
    for k, v in p["org_context"].items():
        if k == "_sources":
            continue
        s = f"  ({src[k]})" if k in src else ""
        print(f"  {k:<36}{v}{s}")

    print("\n멘토님 신호 02 — IP 공백 (역량 보유 · 특허 0건)")
    for g in p["signals"]["ip_gap"]:
        print(f"  ★{g['level']}  {g['name'][:22]:<24}근거: {', '.join(g['evidence'])}")

    print("\n멘토님 신호 03 — 특허 편중")
    pc = p["signals"]["patent_concentration"]
    if pc:
        print(f"  총 {pc['total']}건 / {pc['headcount']}명 · 상위 3명 {pc['top3_share']}%")
        print(f"  특허 0건: {', '.join(pc['zero_patent_staff'])}")

    out = BASE / "organization_profile.json"
    with open(out, "w", encoding="utf-8") as f:
        json.dump(p, f, ensure_ascii=False, indent=2)
    print(f"\n저장: {out.name}")
